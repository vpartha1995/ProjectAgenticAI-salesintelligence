from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Create presentation
prs = Presentation()
slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank slide

# -------------------------
# Title
# -------------------------
title_box = slide.shapes.add_textbox(
    Inches(1), Inches(0.2), Inches(8), Inches(0.8)
)
title_tf = title_box.text_frame
title_tf.clear()

p = title_tf.add_paragraph()
p.text = "Sales Intelligence Agent"
p.font.size = Pt(32)
p.font.bold = True
p.alignment = PP_ALIGN.CENTER

p = title_tf.add_paragraph()
p.text = "Visual System Flow"
p.font.size = Pt(18)
p.alignment = PP_ALIGN.CENTER

# -------------------------
# Box Data (Vertical Flow)
# -------------------------
boxes = [
    ("User\n(Sales Rep)",
     RGBColor(91, 155, 213)),  # Blue

    ("Web Interface (UI)\n• Enter query\n• Select type:\n  Company / News / Lead",
     RGBColor(91, 155, 213)),  # Blue

    ("Tool Selection Layer\n\nCompany → Organizations\nNews → Trends\nLead → Individuals",
     RGBColor(112, 48, 160)),  # Purple

    ("Live Web Search\n(Serper API)\n\n• Google Search\n• Google News",
     RGBColor(0, 176, 240)),  # Teal

    ("AI Summarization\n(Azure OpenAI)\n\n• Key insights\n• Bullet points\n• Clear language",
     RGBColor(112, 48, 160)),  # Purple

    ("Final Output\n\n✔ Summary points\n✔ Source links\n✔ Real-time data",
     RGBColor(0, 176, 80))  # Green
]

# -------------------------
# Draw Boxes
# -------------------------
left = Inches(3)
top = Inches(1.3)
width = Inches(4)
height = Inches(1.0)
gap = Inches(0.2)

for text, color in boxes:
    box = slide.shapes.add_shape(
        autoshape_type_id=1,  # Rounded rectangle
        left=left,
        top=top,
        width=width,
        height=height
    )

    box.fill.solid()
    box.fill.fore_color.rgb = color
    box.fill.fore_color.brightness = 0.4
    box.line.color.rgb = RGBColor(0, 0, 0)

    tf = box.text_frame
    tf.clear()
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(12)
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    top += height + gap

# -------------------------
# Footer (One-line explanation)
# -------------------------
footer = slide.shapes.add_textbox(
    Inches(1), Inches(6.7), Inches(8), Inches(0.6)
)
footer_tf = footer.text_frame
footer_tf.text = (
    "User asks a question → system searches the live web → "
    "AI summarizes → instant sales-ready insights."
)
footer_tf.paragraphs[0].font.size = Pt(14)
footer_tf.paragraphs[0].italic = True
footer_tf.paragraphs[0].alignment = PP_ALIGN.CENTER

# -------------------------
# Save file
# -------------------------
prs.save("Sales_Intelligence_Agent_Visual_Diagram.pptx")

print("PowerPoint created: Sales_Intelligence_Agent_Visual_Diagram.pptx")
